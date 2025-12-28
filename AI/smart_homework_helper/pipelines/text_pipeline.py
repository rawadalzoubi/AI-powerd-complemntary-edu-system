import os
import re
import io
import time
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from pipelines.base_pipeline import BasePipeline

# استيراد المصحح الذكي (الجديد)
from utils.ai_corrector import AICorrector

# استيراد المكتبات الأساسية
try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False

try:
    from PIL import Image
except ImportError:
    print("⚠️ PIL not installed")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

class TextPipeline(BasePipeline):
    def __init__(self, config):
        super().__init__(config)
        print(f"[TextPipeline] Loading Embedding Model: {config.EMBEDDING_MODEL_NAME}...")
        self.embeddings = HuggingFaceEmbeddings(model_name=config.EMBEDDING_MODEL_NAME)
        self.index_path = os.path.join(config.VECTOR_DB_PATH, "text_index")
        
        # ===> إضافة جديدة: تهيئة المصحح <===
        self.ai_helper = AICorrector(api_key=getattr(config, 'GOOGLE_API_KEY', None))

    def run(self, files_to_process=None):
        # (نفس كود run الأصلي تماماً بدون تغيير، فقط تأكد من استدعاء _load_pdf_smart المعدلة)
        # ... [انسخ كود run الذي أرسلته لي كما هو] ...
        # سأختصره هنا لتوفير المساحة، الكود المنطقي لم يتغير، التغيير في الدوال الداخلية
        
        all_docs = []
        is_update = False
        if files_to_process:
            target_files = files_to_process
            is_update = True
        else:
             if not os.path.exists(self.config.DATA_DIR): return
             # البحث في المجلد الرئيسي والمجلدات الفرعية
             target_files = []
             for root, dirs, files in os.walk(self.config.DATA_DIR):
                 for f in files:
                     if f.endswith(('.pdf', '.docx', '.pptx', '.mp4', '.mp3', '.wav')):
                         target_files.append(os.path.join(root, f))

        print(f"[TextPipeline] Processing {len(target_files)} files...")

        for file_path in target_files:
            try:
                ext = os.path.splitext(file_path)[1].lower()
                docs = []
                print(f"  -> Loading: {os.path.basename(file_path)}")
                
                if ext == '.pdf':
                    docs = self._load_pdf_smart(file_path) # <--- هنا التغيير الجوهري
                elif ext == '.docx':
                    loader = Docx2txtLoader(file_path)
                    docs = loader.load()
                elif ext == '.pptx':
                    docs = self._load_pptx(file_path)
                
                if docs:
                    cleaned_docs = self.clean_documents(docs)
                    all_docs.extend(cleaned_docs)
                    print(f"     ✅ Extracted {len(cleaned_docs)} pages.")
            except Exception as e:
                print(f"Error loading {file_path}: {e}")

        if not all_docs: return

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=self.config.CHUNK_SIZE, chunk_overlap=self.config.CHUNK_OVERLAP)
        splits = text_splitter.split_documents(all_docs)

        if is_update and os.path.exists(self.index_path):
             self.vectorstore = FAISS.load_local(self.index_path, self.embeddings, allow_dangerous_deserialization=True)
             self.vectorstore.merge_from(FAISS.from_documents(splits, self.embeddings))
        else:
             self.vectorstore = FAISS.from_documents(splits, self.embeddings)
        
        self.vectorstore.save_local(self.index_path)
        print(f"[TextPipeline] Index saved.")

    def _load_pdf_smart(self, file_path):
        """
        قراءة PDF - نستخدم النص الرقمي أولاً (أسرع وأوفر)
        VLM يُستخدم فقط كـ fallback للصفحات الممسوحة ضوئياً
        """
        if not FITZ_AVAILABLE:
            return PyPDFLoader(file_path).load()
        
        documents = []
        try:
            doc = fitz.open(file_path)
            
            for page_num, page in enumerate(doc):
                text = ""
                
                # 1. أولاً: جرب النص الرقمي (مجاني وسريع)
                digital_text = page.get_text()
                if digital_text and len(digital_text.strip()) > 50:
                    text = digital_text
                    print(f"     📄 Page {page_num+1}: Digital Text ({len(text)} chars)")
                
                # 2. Fallback: VLM فقط إذا النص الرقمي فارغ أو قصير جداً
                elif self.ai_helper.client:
                    try:
                        pix = page.get_pixmap(dpi=200)
                        img_data = pix.tobytes("png")
                        img = Image.open(io.BytesIO(img_data))
                        
                        print(f"     🤖 VLM Scanning Page {page_num+1}...")
                        vlm_text = self.ai_helper.extract_text_from_image(img)
                        
                        if vlm_text and len(vlm_text.strip()) > 20:
                            text = vlm_text
                            print(f"     ✨ VLM Extracted: {len(text)} chars")
                            time.sleep(2)  # تأخير أكبر لتجنب rate limit
                    except Exception as e:
                        print(f"     ⚠️ VLM Skipped for page {page_num+1}: {e}")

                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={"source": os.path.basename(file_path), "page": page_num + 1}
                    ))
            
            return documents
            
        except Exception as e:
            print(f"     ❌ Error processing PDF: {e}")
            return PyPDFLoader(file_path).load()

    def _load_pptx(self, file_path):
        if not Presentation: return []
        try:
            prs = Presentation(file_path)
            text_content = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if slide_text:
                    doc = Document(page_content="\n".join(slide_text), metadata={"source": os.path.basename(file_path), "page": i + 1})
                    text_content.append(doc)
            return text_content
        except: return []

    def clean_documents(self, docs):
        cleaned = []
        for doc in docs:
            text = doc.page_content
            text = re.sub(r'\n{3,}', '\n\n', text)
            text = re.sub(r'[ \t]+', ' ', text).strip()
            doc.page_content = text
            if len(text) > 10: cleaned.append(doc)
        return cleaned
    
    def load_index(self):
        if os.path.exists(self.index_path):
            self.vectorstore = FAISS.load_local(self.index_path, self.embeddings, allow_dangerous_deserialization=True)
            self.retriever = self.vectorstore.as_retriever(search_kwargs={"k": self.config.TOP_K_RETRIEVAL})


    def search(self, query, k=None):
        """Search the index for relevant documents."""
        if not self.vectorstore:
            self.load_index()
        if not self.vectorstore:
            return []
        search_k = k or self.config.TOP_K_RETRIEVAL
        return self.vectorstore.similarity_search(query, k=search_k)
