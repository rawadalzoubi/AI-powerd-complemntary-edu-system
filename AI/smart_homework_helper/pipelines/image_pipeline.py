import os
import fitz  # PyMuPDF
import io
import numpy as np
import faiss
import zipfile
import pickle
from PIL import Image
from sentence_transformers import SentenceTransformer
from pipelines.base_pipeline import BasePipeline

# محاولة استيراد pptx
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

class ImagePipeline(BasePipeline):
    def __init__(self, config):
        super().__init__(config)
        print(f"[ImagePipeline] Loading Embedding Model: {config.IMAGE_EMBEDDING_MODEL_NAME}...")
        self.model = SentenceTransformer(config.IMAGE_EMBEDDING_MODEL_NAME)
        
        self.index_path = os.path.join(config.VECTOR_DB_PATH, "image_index")
        self.metadata = []
        self.index = None

    def run(self, files_to_process=None):
        # 1. إعداد المجلدات
        output_dir = os.path.join(self.config.DATA_DIR, "extracted_images")
        if not os.path.exists(output_dir): os.makedirs(output_dir)
        if not os.path.exists(self.config.VECTOR_DB_PATH): os.makedirs(self.config.VECTOR_DB_PATH)

        # 2. إعداد الفهرس
        test_vec = self.model.encode(Image.new('RGB', (50, 50)))
        dimension = len(test_vec)
        self.index = faiss.IndexFlatL2(dimension)
        self.metadata = []

        # 3. تحديد الملفات
        supported_ext = ('.pdf', '.pptx', '.docx')
        
        if files_to_process:
            target_files = [f for f in files_to_process if f.lower().endswith(supported_ext)]
        else:
            if not os.path.exists(self.config.DATA_DIR): return
            # البحث في المجلد الرئيسي والمجلدات الفرعية
            target_files = []
            for root, dirs, files in os.walk(self.config.DATA_DIR):
                # تجاهل مجلدات الإخراج
                if 'extracted_images' in root or 'debug_extracted_texts' in root:
                    continue
                for f in files:
                    if f.lower().endswith(supported_ext):
                        target_files.append(os.path.join(root, f))

        if not target_files:
            print("[ImagePipeline] No supported files to process.")
            return

        print(f"[ImagePipeline] Processing {len(target_files)} files (Full Page Indexing)...")
        
        for file_path in target_files:
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.pdf':
                self._process_pdf_pages(file_path, output_dir)
            elif ext == '.pptx':
                self._process_pptx_slides(file_path, output_dir)
            elif ext == '.docx':
                self._process_docx_images(file_path, output_dir)
            
        # 4. الحفظ
        if self.metadata:
            print(f"[ImagePipeline] Saving index with {len(self.metadata)} pages/images...")
            faiss.write_index(self.index, self.index_path + ".faiss")
            with open(self.index_path + "_meta.pkl", "wb") as f:
                pickle.dump(self.metadata, f)
            print(f"[ImagePipeline] Index saved.")
        else:
            print("[ImagePipeline] No images extracted.")

    # --- معالجة PDF (تحويل كل صفحة لصورة) ---
    def _process_pdf_pages(self, pdf_path, output_dir):
        try:
            filename = os.path.basename(pdf_path)
            doc = fitz.open(pdf_path)
            print(f"  -> Converting PDF pages to images: {filename}...")

            for page_index, page in enumerate(doc):
                try:
                    # تحويل الصفحة كاملة لصورة عالية الدقة (dpi=150 كافية للفهرسة)
                    pix = page.get_pixmap(dpi=150)
                    img_data = pix.tobytes("png")
                    
                    # استخراج نص الصفحة ليكون سياقاً
                    page_text = page.get_text()
                    
                    # حفظ وفهرسة الصفحة كصورة
                    self._save_bytes_image(
                        img_data, 
                        output_dir, 
                        filename, 
                        page_index, 
                        page_text, 
                        "full_page" # نميزها بأنها صفحة كاملة
                    )
                    
                except Exception as e:
                    print(f"     Failed to process page {page_index+1}: {e}")

            print(f"     ✅ Processed {len(doc)} pages.")

        except Exception as e:
            print(f"Error processing PDF {filename}: {e}")

    # --- معالجة PPTX (تحويل كل شريحة لصورة - إن أمكن، أو استخراج الصور) ---
    # ملاحظة: python-pptx لا تدعم تحويل الشريحة لصورة مباشرة بسهولة بدون مكتبات نظام.
    # لذا سنستمر في استخراج الصور من الشرائح، ولكن سنسميها باسم الشريحة.
    def _process_pptx_slides(self, pptx_path, output_dir):
        if not Presentation: return
        try:
            filename = os.path.basename(pptx_path)
            print(f"  -> Scanning PPTX: {filename}...")
            prs = Presentation(pptx_path)

            for i, slide in enumerate(prs.slides):
                # نجمع نص الشريحة
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                slide_text = " ".join(texts)
                
                # نستخرج الصور داخل الشريحة
                for shape in slide.shapes:
                    if hasattr(shape, "shape_type") and shape.shape_type == 13: # PICTURE
                        try:
                            image = shape.image
                            self._save_bytes_image(image.blob, output_dir, filename, i, slide_text, f"slide_img_{shape.shape_id}")
                        except: pass
        except: pass

    # --- معالجة DOCX (استخراج الصور) ---
    def _process_docx_images(self, docx_path, output_dir):
        try:
            filename = os.path.basename(docx_path)
            print(f"  -> Scanning DOCX: {filename}...")
            with zipfile.ZipFile(docx_path) as z:
                media_files = [f for f in z.namelist() if f.startswith('word/media/')]
                for media_file in media_files:
                    self._save_bytes_image(z.read(media_file), output_dir, filename, 0, f"Image from {filename}", os.path.basename(media_file))
        except: pass

    # --- دوال مساعدة للحفظ والفهرسة ---
    def _save_bytes_image(self, image_bytes, output_dir, source, page_idx, context, img_name):
        try:
            pil_img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            
            # حفظ الملف
            # نستخدم اسم مميز: الملف_الصفحة_النوع.png
            save_filename = f"{source}_p{page_idx+1}_{img_name}.png"
            save_path = os.path.join(output_dir, save_filename)
            pil_img.save(save_path)

            # التضمين (CLIP)
            vector = self.model.encode(pil_img)
            self.index.add(np.array([vector]).astype('float32'))
            
            # الميتاداتا
            self.metadata.append({
                "image_path": save_path,
                "source": source,
                "page_number": page_idx + 1,
                "context_text": context[:1000] + "...", # نحفظ نصاً أطول للسياق
                "type": img_name
            })
        except Exception as e:
            pass

    def load_index(self):
        try:
            if os.path.exists(self.index_path + ".faiss"):
                print("[ImagePipeline] Loading index...")
                self.index = faiss.read_index(self.index_path + ".faiss")
                with open(self.index_path + "_meta.pkl", "rb") as f: self.metadata = pickle.load(f)
            else: print("[ImagePipeline] No index found.")
        except: pass

    def search(self, query_image_file):
        """
        البحث عن صور مشابهة.
        نرجع أفضل النتائج مع درجة الثقة.
        """
        if not self.index or self.index.ntotal == 0: 
            return []
        try:
            img = Image.open(query_image_file).convert("RGB")
            vec = self.model.encode(img)
            distances, indices = self.index.search(np.array([vec]).astype('float32'), 3)
            
            results = []
            print(f"🔍 Image search distances: {distances[0]}")
            
            for i, idx in enumerate(indices[0]):
                if idx != -1 and idx < len(self.metadata):
                    distance = float(distances[0][i])
                    result = self.metadata[idx].copy()
                    
                    # حساب مستوى الثقة
                    if distance < 1.0:
                        confidence = "عالية جداً ✅"
                    elif distance < 2.0:
                        confidence = "عالية"
                    elif distance < 3.0:
                        confidence = "متوسطة"
                    else:
                        confidence = "منخفضة ⚠️"
                    
                    result['confidence'] = confidence
                    result['distance'] = distance
                    results.append(result)
                    print(f"   📄 {result['source']} p.{result['page_number']} (distance: {distance:.2f}, {confidence})")
            
            return results
        except Exception as e:
            print(f"❌ Image search error: {e}")
            return []